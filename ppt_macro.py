import os, sys

from pptx import Presentation 
from pptx.util import Inches, Pt

from glob import glob
from PIL import Image
from pdf2image import convert_from_path
from io import BytesIO
import re

####### CONFIG ###########
OUTPUT_NAME = 'test'
PATH_TO_OUTPUT = './slides/'

path_to_img = './latest_plots/*.pdf'
img_rows = 2 # Rows of images per slide
img_cols = 2 # Col  of images per slide

imgs_on_slide    = img_rows*img_cols
num_of_keni_type = 12    # Only have images of the same type slide
num_of_data_type = 0
sel_cuts = {
    'nb0':'nBottom=0',        'nbg0':'nBottom>0', 'nbg1':'nBottom>1',
    'nRt0':'nResolvedTops=0', 'nRtg0':'nResolvedTops>0', 'nRtg1':'nResolvedTops>1',
    'nMt0':'nMergedTops=0',   'nMtg0':'nMergedTops>0',   'nMtg1':'nMergedTops>1',
    'nRMtl3':'nResolvedTops_drLeptonCleaned+nMergedTops_drLeptonCleaned)<3'}

image_files = list(glob(path_to_img))
image_files.sort()
###########################
prs = Presentation()
# TITLE SLIDE
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Baylor Meeting"
slide.placeholders[1].text = 'Bryan Caraway'
##
# CUT DISCRIPTION SLIDE
text_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(text_slide_layout)
body_shape = slide.shapes.placeholders[1]
p = body_shape.text_frame.paragraphs[0]
p.text = 'Variable cuts: '
p = body_shape.text_frame.add_paragraph()
text_slide_string = ''
for cut in sel_cuts.values():
    text_slide_string = text_slide_string + cut + '\n'
p.text = text_slide_string
p.level = 1
##
#TABLE SLIDES
table_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(table_slide_layout)
rows = cols = 4
left   = Inches(0.5)
top    = Inches(1.5)
width  = Inches(9.0)
height = Inches(4.0)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
table.cell(0,0).text = 'nBottom\nZpt>300'
for i in range(1,4):
    table.cell(0,i).text = 'nResolvedTops'
    table.cell(i,0).text = 'nMergedTops'
##
# IMAGE SLIDES
blank_slide_layout = prs.slide_layouts[6]
count1 = count2 = 0 # I think I need this  
for i,image_file in enumerate(image_files):
    print (image_file)
    slideimgs = convert_from_path(image_file, 100, fmt='ppm', thread_count=1, use_cropbox=True)
    
    for slideimg in slideimgs:
        imagefile = BytesIO()
        slideimg.save(imagefile, format='tiff')
        imagefile.seek(0)
        #width, height = slideimg.size

        # Add slide ##
        if ( count1 % imgs_on_slide == 0 or (count2 % num_of_keni_type == 0 and count2 != 0) or ("Data" in image_file and i % num_of_data_type == 0)):
            slide = prs.slides.add_slide(blank_slide_layout)
            titleBox = slide.shapes.add_textbox(Inches(2), Inches(0.0),
                                                    prs.slide_width*.50, Inches(.8))
            p = titleBox.text_frame.paragraphs[0]
            title_text = ''
            if ('mu'   in image_file):
                title_text = 'Muon'
            if ('elec'   in image_file):
                title_text = 'Electron'
            if ('NoZMassCut' in image_file):
                title_text = title_text + " --- " + "No Z Mass Cut"
            p.text = title_text
            p.font.size = Pt(30)
            count1 = 0
        ##
        i_row = Inches(count1 % img_rows)
        i_col = Inches(int(count1 % imgs_on_slide >= img_cols))
        left  = Inches(1)+i_row*4.5
        top   = Inches(.95)+i_col*3.3
        width = prs.slide_width*.32
        # Add Picture to slide with description text
        pic      = slide.shapes.add_picture(imagefile, left, top,
                                            width=width)
        add_text = slide.shapes.add_textbox(left+width*.30, top+pic.height*.2,
                                            prs.slide_width*.20, Inches(2))
        p        = add_text.text_frame.paragraphs[0]
        plot_sels = image_file.split('_')
        plot_text = ''
        for sel in plot_sels:
            #if re.sub(r'\d+', '',sel) in sel_cuts:
            if sel in sel_cuts:
                plot_text = plot_text+sel+'\n'
        p.text = plot_text
        ##
        if ("Data" not in image_file):
            count2 += 1 
        count1 += 1
    
##
prs.save( PATH_TO_OUTPUT+OUTPUT_NAME+'.pptx')



    
    
    
